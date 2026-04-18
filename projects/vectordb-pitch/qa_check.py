# -*- coding: utf-8 -*-
from pptx import Presentation
import sys

prs = Presentation("shared-agent-intelligence.pptx")

print("=" * 80)
print("POWERPOINT QA CHECK REPORT")
print("=" * 80)
print("Total slides: {}".format(len(prs.slides)))
print()

SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 5.625
TITLE_MIN, TITLE_MAX = 28, 44
BODY_MIN, BODY_MAX = 10, 16

issues_found = False

for slide_idx, slide in enumerate(prs.slides, 1):
    print("\n" + "="*80)
    print("SLIDE {}".format(slide_idx))
    print("="*80)
    
    slide_issues = []
    
    if len(slide.shapes) == 0:
        print("  [WARNING] No shapes")
        continue
    
    print("  Shapes: {}".format(len(slide.shapes)))
    
    for shape_idx, shape in enumerate(slide.shapes, 1):
        shape_name = shape.name
        
        if not hasattr(shape, 'text_frame'):
            continue
        
        text_frame = shape.text_frame
        try:
            shape_text = text_frame.text.strip()
        except:
            shape_text = ""
        
        x_inches = shape.left.inches
        y_inches = shape.top.inches
        w_inches = shape.width.inches
        h_inches = shape.height.inches
        
        print("\n  [{}] {}".format(shape_idx, shape_name))
        print("      Pos: ({:.2f}\", {:.2f}\")".format(x_inches, y_inches))
        print("      Size: {:.2f}\" x {:.2f}\"".format(w_inches, h_inches))
        
        if not shape_text:
            print("      [EMPTY]")
            slide_issues.append("- {}: Empty text".format(shape_name))
            issues_found = True
        else:
            try:
                preview = shape_text.encode('ascii', 'ignore').decode('ascii')[:50]
            except:
                preview = "[unicode text]"
            print("      Text: {}".format(preview))
        
        x_end = x_inches + w_inches
        y_end = y_inches + h_inches
        
        if x_end > SLIDE_WIDTH:
            print("      [OVERFLOW-X] {:.2f}\" > {:.2f}\"".format(x_end, SLIDE_WIDTH))
            slide_issues.append("- {}: X overflow".format(shape_name))
            issues_found = True
        
        if y_end > SLIDE_HEIGHT:
            print("      [OVERFLOW-Y] {:.2f}\" > {:.2f}\"".format(y_end, SLIDE_HEIGHT))
            slide_issues.append("- {}: Y overflow".format(shape_name))
            issues_found = True
        
        if shape_text:
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                for run in paragraph.runs:
                    if run.font.size:
                        font_pt = run.font.size.pt
                        is_title = para_idx == 0 and font_pt > 20
                        
                        if is_title:
                            if not (TITLE_MIN <= font_pt <= TITLE_MAX):
                                print("      Font: {}pt [TITLE OUT OF RANGE {}-{}]".format(font_pt, TITLE_MIN, TITLE_MAX))
                                slide_issues.append("- {}: Title {}pt".format(shape_name, font_pt))
                                issues_found = True
                        else:
                            if not (BODY_MIN <= font_pt <= BODY_MAX):
                                print("      Font: {}pt [BODY OUT OF RANGE {}-{}]".format(font_pt, BODY_MIN, BODY_MAX))
                                slide_issues.append("- {}: Body {}pt".format(shape_name, font_pt))
                                issues_found = True
    
    text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
    for i in range(len(text_shapes)):
        for j in range(i+1, len(text_shapes)):
            s1, s2 = text_shapes[i], text_shapes[j]
            x1, y1 = s1.left.inches, s1.top.inches
            w1, h1 = s1.width.inches, s1.height.inches
            x2, y2 = s2.left.inches, s2.top.inches
            w2, h2 = s2.width.inches, s2.height.inches
            
            ox = x1 < x2 + w2 and x1 + w1 > x2
            oy = y1 < y2 + h2 and y1 + h1 > y2
            
            if ox and oy:
                t1 = s1.text_frame.text.strip()
                t2 = s2.text_frame.text.strip()
                if t1 or t2:
                    print("\n  [OVERLAP] '{}' and '{}'".format(s1.name, s2.name))
                    slide_issues.append("- Overlap: {} + {}".format(s1.name, s2.name))
                    issues_found = True
    
    if slide_issues:
        print("\n  ISSUES:")
        for issue in slide_issues:
            print("  " + issue)
    else:
        print("\n  [PASS]")

print("\n" + "="*80)
if issues_found:
    print("SUMMARY: Issues found")
    sys.exit(1)
else:
    print("SUMMARY: All checks passed!")
    sys.exit(0)
